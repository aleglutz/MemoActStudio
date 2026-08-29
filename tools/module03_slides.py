"""Build the Module 03 deck from the run sheet, so there is only one source.

`projects/module03/Slides for Module 03.md` is authoritative: it carries the
wording, the order and the minute budget, and it is what gets edited. This
script renders it as a .pptx, which is the thing that goes on the screen -- via
PowerPoint, LibreOffice, or an import into Google Slides.

Nothing is written into the deck by hand. Editing a slide means editing the
markdown and running this again; a deck that had been touched afterwards would
be a second source of truth, and the run sheet would start losing arguments to
it.

Three conventions in the markdown decide what a block becomes:

    > **Say:** ...            speaker notes -- the presenter's line, not a slide
    > **anything else**       a pull quote, on the slide, because it is the point
    *(9 min for S16-S18.)*    a minute budget: into the notes and the corner

and a `stills/...png` or `out/...mp4` named in backticks becomes the picture.
For a video that means one frame lifted with ffmpeg at a moment chosen below --
the deck cannot play the clip, and a still of it is an honest placeholder for
the file the presenter opens beside the deck.

Usage:
    python tools/module03_slides.py
    python tools/module03_slides.py --out somewhere/else.pptx --no-media
"""

import argparse
import os
import re
import shutil
import subprocess
import sys
import tempfile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PROJECT = os.path.join(HERE, "projects", "module03")
SOURCE = os.path.join(PROJECT, "Slides for Module 03.md")
DEFAULT_OUT = os.path.join(PROJECT, "out", "Module_03.pptx")

# Where to lift a poster frame from each clip. Chosen for what is on screen at
# that second, not for the arithmetic middle: L1 wants a caption up (the third
# quote runs 18.0-25.5 in quotes.ass), L2 wants the press crowding the table,
# L3 wants both halves of the A/B busy.
POSTER_AT = {
    "out/L1_quotes.mp4": 20.0,
    "out/L2_sound.mp4": 15.0,
    "out/L3_ab_zoom_30s.mp4": 15.0,
    "out/L3_split_30s.mp4": 15.0,
    "out/L4_ab.mp4": 3.0,
    # The reel: 1:01 is the archival shot itself, which is what S4 is about.
    "../legends_of_surrender/out/reel_with_hook.mp4": 61.0,
    # The source film: 4:20, inside the stretch every level below works on.
    "../legends_of_surrender/sources/videos/MBK_KAPFILM_FINAL.mp4": 260.0,
}

# Warm near-black rather than pure black: the footage is a warm-tinted scan and
# a neutral ground makes it look like a fault in the transfer.
INK = RGBColor(0x14, 0x11, 0x0E)
PAPER = RGBColor(0xF2, 0xEC, 0xE4)
ACCENT = RGBColor(0xD1, 0x82, 0x4A)
MUTED = RGBColor(0x94, 0x89, 0x7C)
RULE = RGBColor(0x33, 0x2C, 0x26)
PANEL = RGBColor(0x1E, 0x1A, 0x16)
STRIPE = RGBColor(0x1A, 0x16, 0x13)

TITLE_FONT = "Georgia"
BODY_FONT = "Arial"
MONO_FONT = "Consolas"

W, H = 13.333, 7.5
MARGIN = 0.85
BODY_TOP = 2.15
BODY_BOTTOM = 6.85
FULL_WIDTH = W - 2 * MARGIN
SPLIT_WIDTH = 6.55  # text column when a picture takes the right-hand side


# --------------------------------------------------------------------------
# reading the run sheet
# --------------------------------------------------------------------------

def parse(md):
    """Split the markdown into slides, and keep the preamble for the run sheet.

    The run-sheet table sits above S1 and is the only table outside a slide,
    which is what identifies it.
    """
    preamble, slides = [], []
    current = None
    for line in md.splitlines():
        head = re.match(r"^## S(\d+)\s+(.*)$", line)
        if head:
            current = {"n": int(head.group(1)), "title": head.group(2).strip(), "lines": []}
            slides.append(current)
        elif current is None:
            preamble.append(line)
        else:
            current["lines"].append(line)
    for slide in slides:
        slide["raw"] = "\n".join(slide["lines"])
        slide["blocks"] = blocks_of(slide["lines"])
    return preamble, slides


def is_separator_row(cells):
    return bool(cells) and all(re.fullmatch(r":?-{2,}:?", c) for c in cells if c)


def table_row(line):
    return [c.strip() for c in line.strip().strip("|").split("|")]


def blocks_of(lines):
    """Turn one slide's lines into typed blocks, in order."""
    out = []
    i, n = 0, len(lines)
    while i < n:
        line = lines[i]
        if not line.strip() or line.strip() == "---":
            i += 1
            continue

        if line.startswith("|"):
            rows = []
            while i < n and lines[i].startswith("|"):
                cells = table_row(lines[i])
                if not is_separator_row(cells):
                    rows.append(cells)
                i += 1
            out.append(("table", rows))
            continue

        if line.startswith(">"):
            quote = []
            while i < n and lines[i].startswith(">"):
                quote.append(lines[i].lstrip(">").strip())
                i += 1
            out.append(("quote", join_wrapped(quote)))
            continue

        if line.startswith("### "):
            out.append(("h3", line[4:].strip()))
            i += 1
            continue

        if re.match(r"^(-|\d+\.)\s+", line):
            ordered = not line.startswith("-")
            items = []
            while i < n:
                item = re.match(r"^(?:-|\d+\.)\s+(.*)$", lines[i])
                if item:
                    items.append(item.group(1).strip())
                    i += 1
                elif items and lines[i].startswith(("   ", "\t")):
                    items[-1] += " " + lines[i].strip()
                    i += 1
                else:
                    break
            out.append(("numbers" if ordered else "bullets", items))
            continue

        if re.fullmatch(r"\*\(.*\)\*", line.strip()):
            out.append(("timing", line.strip()[2:-2]))
            i += 1
            continue

        para = []
        while i < n and lines[i].strip() \
                and not lines[i].startswith(("|", ">", "-", "#")) \
                and not re.match(r"^\d+\.\s", lines[i]) \
                and not re.fullmatch(r"\*\(.*\)\*", lines[i].strip()):
            para.append(lines[i].strip())
            i += 1
        out.append(("para", " ".join(para)))
    return out


def join_wrapped(quote_lines):
    """Rejoin a blockquote's hard-wrapped lines, keeping its blank-line breaks."""
    paras, current = [], []
    for line in quote_lines:
        if line:
            current.append(line)
        elif current:
            paras.append(" ".join(current))
            current = []
    if current:
        paras.append(" ".join(current))
    return paras


def is_note(quote_paras):
    """A blockquote is the presenter's line, not a slide, when it opens like one.

    `> **Say:**` is the obvious case. The two others are the same thing wearing
    a different hat -- a correction to make aloud, and the cue for handing the
    room back from the designer. Everything else in a blockquote stays on the
    slide, and those are the loudest thing on it by design.
    """
    if not quote_paras:
        return False
    return quote_paras[0].startswith(("**Say", "**Hand-back cue", "**Two corrections"))


def is_callout(quote_paras):
    """A warning the room has to see while it works, not one it hears once."""
    return bool(quote_paras) and quote_paras[0].startswith("**The trap")


# --------------------------------------------------------------------------
# inline markdown
# --------------------------------------------------------------------------

INLINE = re.compile(r"(\*\*.+?\*\*|\*[^*]+?\*|`[^`]+?`)")


def write_runs(paragraph, text, size, color=PAPER, bold_color=None, base_bold=False):
    """Fill a paragraph with runs, honouring **bold**, *italic* and `code`."""
    for piece in INLINE.split(text):
        if not piece:
            continue
        run = paragraph.add_run()
        font = run.font
        font.size = Pt(size)
        font.name = BODY_FONT
        font.color.rgb = color
        font.bold = base_bold
        if piece.startswith("**") and piece.endswith("**"):
            run.text = piece[2:-2]
            font.bold = True
            font.color.rgb = bold_color or color
        elif piece.startswith("`") and piece.endswith("`"):
            run.text = piece[1:-1]
            font.name = MONO_FONT
            font.size = Pt(size - 1)
            font.color.rgb = ACCENT
        elif piece.startswith("*") and piece.endswith("*"):
            run.text = piece[1:-1]
            font.italic = True
        else:
            run.text = piece


def plain(text):
    return re.sub(r"[*`]", "", text)


# --------------------------------------------------------------------------
# measuring, so a long slide shrinks instead of overflowing
# --------------------------------------------------------------------------

CHARS_PER_INCH = 5.6  # 18 pt Arial, near enough to count lines with


def wrapped_lines(text, width_in, size):
    per_line = max(8, int(width_in * CHARS_PER_INCH * 18 / size))
    return max(1, -(-len(plain(text)) // per_line))


def body_size(blocks, width):
    """One body size for the slide, chosen from how much has to fit on it."""
    lines = 0
    for kind, payload in blocks:
        if kind in ("para", "h3"):
            lines += wrapped_lines(payload, width, 18) + 1
        elif kind == "quote":
            lines += sum(wrapped_lines(t, width, 18) for t in payload) + 1
        elif kind in ("bullets", "numbers"):
            lines += sum(wrapped_lines(t, width - 0.4, 18) for t in payload) + 1
        elif kind == "table":
            lines += len(payload) * 1.4 + 1
    if lines <= 11:
        return 18
    if lines <= 15:
        return 16
    if lines <= 20:
        return 14
    return 12.5


# --------------------------------------------------------------------------
# drawing
# --------------------------------------------------------------------------

def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def ground(slide, prs):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                   prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = INK
    shape.line.fill.background()
    shape.shadow.inherit = False


CHAPTERS = [
    (1, 4, "Opening"),
    (5, 5, "Moodboard"),
    (6, 7, "The film"),
    (8, 8, "The ladder"),
    (9, 10, "Level 1 — quotes"),
    (11, 11, "Break"),
    (12, 15, "Level 2 — sound"),
    (16, 18, "Level 3 — upscale"),
    (19, 21, "Level 4 — colour"),
    (22, 23, "Close"),
]


def chapter_of(n):
    """The block a slide belongs to, printed small above its title."""
    for lo, hi, name in CHAPTERS:
        if lo <= n <= hi:
            return name
    return "Presenter"


def header(slide, n, timing):
    frame = textbox(slide, MARGIN, 0.52, FULL_WIDTH, 0.32)
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = "S%d   ·   %s" % (n, chapter_of(n).upper())
    run.font.size = Pt(11)
    run.font.name = BODY_FONT
    run.font.bold = True
    run.font.color.rgb = ACCENT
    minutes = re.search(r"(\d+)\s*min", timing or "")
    if minutes:
        tail = p.add_run()
        tail.text = "   ·   %s MIN" % minutes.group(1)
        tail.font.size = Pt(11)
        tail.font.name = BODY_FONT
        tail.font.color.rgb = MUTED
    rect(slide, MARGIN, 0.92, FULL_WIDTH, 0.02, RULE)


def title(slide, text, width):
    frame = textbox(slide, MARGIN, 1.10, width, 0.95)
    p = frame.paragraphs[0]
    p.line_spacing = 1.06
    write_runs(p, text, 30 if len(plain(text)) < 62 else 24, PAPER, bold_color=PAPER)
    for run in p.runs:
        run.font.name = TITLE_FONT
        run.font.bold = True


def footer(slide, n, total):
    frame = textbox(slide, MARGIN, H - 0.5, FULL_WIDTH, 0.28)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Module 03 · Intervening in archival footage · %d / %d" % (n, total)
    run.font.size = Pt(9)
    run.font.name = BODY_FONT
    run.font.color.rgb = MUTED


def draw_para(slide, y, text, width, size, muted=False):
    height = wrapped_lines(text, width, size) * size / 72.0 * 1.35
    frame = textbox(slide, MARGIN, y, width, height)
    p = frame.paragraphs[0]
    p.line_spacing = 1.3
    write_runs(p, text, size, MUTED if muted else PAPER, bold_color=PAPER)
    return y + height + 0.16


def draw_list(slide, y, items, width, size, ordered):
    for idx, item in enumerate(items, 1):
        mark = textbox(slide, MARGIN, y, 0.4, 0.3)
        run = mark.paragraphs[0].add_run()
        run.text = "%d." % idx if ordered else "—"
        run.font.size = Pt(size)
        run.font.name = BODY_FONT
        run.font.bold = ordered
        run.font.color.rgb = ACCENT

        height = wrapped_lines(item, width - 0.45, size) * size / 72.0 * 1.35
        frame = textbox(slide, MARGIN + 0.45, y, width - 0.45, height)
        p = frame.paragraphs[0]
        p.line_spacing = 1.3
        write_runs(p, item, size, PAPER, bold_color=PAPER)
        y += height + 0.13
    return y + 0.06


def draw_quote(slide, y, paras, width, size):
    text_size = min(21, size + 3)
    inner = width - 0.6
    height = sum(wrapped_lines(t, inner, text_size) for t in paras)
    height = height * text_size / 72.0 * 1.42 + 0.30
    rect(slide, MARGIN, y, 0.055, height, ACCENT)
    frame = textbox(slide, MARGIN + 0.42, y + 0.13, inner, height - 0.22)
    for idx, para in enumerate(paras):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.line_spacing = 1.28
        p.space_before = Pt(0 if idx == 0 else 8)
        write_runs(p, para, text_size, PAPER, bold_color=ACCENT)
        for run in p.runs:
            run.font.name = TITLE_FONT
    return y + height + 0.20


def draw_callout(slide, y, paras, width, size):
    text_size = min(15, size)
    inner = width - 0.85
    height = sum(wrapped_lines(t, inner, text_size) for t in paras)
    height = height * text_size / 72.0 * 1.40 + 0.50
    rect(slide, MARGIN, y, width, height, PANEL)
    rect(slide, MARGIN, y, 0.055, height, ACCENT)
    frame = textbox(slide, MARGIN + 0.40, y + 0.24, inner, height - 0.44)
    for idx, para in enumerate(paras):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.line_spacing = 1.28
        p.space_before = Pt(0 if idx == 0 else 7)
        write_runs(p, para, text_size, PAPER, bold_color=ACCENT)
    return y + height + 0.20


def draw_table(slide, y, rows, width, size):
    cols = max(len(r) for r in rows)
    rows = [r + [""] * (cols - len(r)) for r in rows]

    # Column width from the longest cell, damped: a prompt column six times the
    # width of a seed column would leave the seeds unreadable.
    weights = [max(len(plain(r[c])) for r in rows) ** 0.72 for c in range(cols)]
    widths = [max(0.75, width * w / sum(weights)) for w in weights]
    scale = width / sum(widths)
    widths = [w * scale for w in widths]

    cell_size = min(13, size)
    heights = [
        max(wrapped_lines(row[c], widths[c] - 0.24, cell_size) for c in range(cols))
        * cell_size / 72.0 * 1.45 + 0.16
        for row in rows
    ]

    table = slide.shapes.add_table(len(rows), cols, Inches(MARGIN), Inches(y),
                                   Inches(width), Inches(sum(heights))).table
    table.first_row = True
    table.horz_banding = False
    for c in range(cols):
        table.columns[c].width = Inches(widths[c])
    for r in range(len(rows)):
        table.rows[r].height = Inches(heights[r])

    for r, row in enumerate(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.11)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if r == 0 else (INK if r % 2 else STRIPE)
            p = cell.text_frame.paragraphs[0]
            p.line_spacing = 1.2
            write_runs(p, row[c], cell_size, ACCENT if r == 0 else PAPER,
                       bold_color=PAPER, base_bold=(r == 0))
    return y + sum(heights) + 0.22


def draw_cue(slide, y, path, width):
    frame = textbox(slide, MARGIN, y, width, 0.32)
    p = frame.paragraphs[0]
    lead = p.add_run()
    lead.text = "▶  play   "
    lead.font.size = Pt(13)
    lead.font.name = BODY_FONT
    lead.font.bold = True
    lead.font.color.rgb = ACCENT
    name = p.add_run()
    name.text = path
    name.font.size = Pt(13)
    name.font.name = MONO_FONT
    name.font.color.rgb = PAPER
    return y + 0.40


# --------------------------------------------------------------------------
# pictures
# --------------------------------------------------------------------------

def media_for(raw):
    """The picture a slide asks for: a still if it names one, else a clip."""
    paths = re.findall(
        r"`((?:\.\./[\w-]+/)?(?:stills|out|sources/videos)/[^`]+)`", raw)
    stills = [p for p in paths if p.endswith(".png")]
    if stills:
        return stills[0], None
    clips = [p for p in paths if p.endswith(".mp4")]
    return None, clips[0] if clips else None


def poster(clip, tmp):
    """One frame of a clip, so the slide shows what the presenter will play."""
    src = os.path.join(PROJECT, clip)
    if not os.path.exists(src):
        print("  no clip at %s -- slide has no picture" % clip)
        return None
    dst = os.path.join(tmp, os.path.basename(clip)[:-4] + ".png")
    try:
        subprocess.run(
            ["ffmpeg", "-v", "error", "-y", "-ss", str(POSTER_AT.get(clip, 5.0)),
             "-i", src, "-frames:v", "1", dst],
            check=True)
    except (OSError, subprocess.CalledProcessError) as exc:
        print("  poster frame failed for %s (%s)" % (clip, exc))
        return None
    return dst


def aspect_of(path):
    with Image.open(path) as img:
        return img.width / img.height


def place_picture(slide, path, band, top):
    """Full width under the text if the picture is a strip, else a right column."""
    aspect = aspect_of(path)
    if band:
        w = FULL_WIDTH
        h = w / aspect
        if top + h > BODY_BOTTOM:
            h = BODY_BOTTOM - top
            w = h * aspect
        slide.shapes.add_picture(path, Inches(MARGIN + (FULL_WIDTH - w) / 2),
                                 Inches(top), Inches(w), Inches(h))
        return
    w = W - MARGIN - (MARGIN + SPLIT_WIDTH + 0.55)
    h = w / aspect
    if h > BODY_BOTTOM - BODY_TOP:
        h = BODY_BOTTOM - BODY_TOP
        w = h * aspect
    slide.shapes.add_picture(path, Inches(W - MARGIN - w), Inches(BODY_TOP),
                             Inches(w), Inches(h))


# --------------------------------------------------------------------------
# the deck
# --------------------------------------------------------------------------

def title_slide(prs, deck_title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ground(slide, prs)
    rect(slide, MARGIN, 2.55, 1.5, 0.06, ACCENT)
    frame = textbox(slide, MARGIN, 2.95, 10.6, 2.4)
    p = frame.paragraphs[0]
    p.line_spacing = 1.12
    run = p.add_run()
    run.text = deck_title
    run.font.size = Pt(38)
    run.font.name = TITLE_FONT
    run.font.bold = True
    run.font.color.rgb = PAPER

    sub = textbox(slide, MARGIN, 5.85, 10.6, 0.9)
    for idx, line in enumerate(subtitle):
        p = sub.paragraphs[0] if idx == 0 else sub.add_paragraph()
        p.space_before = Pt(0 if idx == 0 else 5)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.name = BODY_FONT
        run.font.color.rgb = ACCENT if idx == 0 else MUTED
    return slide


def notes_of(data):
    """Everything the presenter needs and the room must not read off a wall."""
    out = []
    for kind, payload in data["blocks"]:
        if kind == "quote" and is_note(payload):
            out.extend(plain(p) for p in payload)
        elif kind == "timing":
            out.append("[ %s ]" % plain(payload))
    return "\n\n".join(out)


def build(md_path, out_path, with_media=True):
    with open(md_path, encoding="utf-8") as fh:
        preamble, slides = parse(fh.read())

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    total = len(slides)
    tmp = tempfile.mkdtemp(prefix="module03_slides_")
    try:
        for data in slides:
            render(prs, data, total, tmp, with_media)
        run_sheet_slide(prs, preamble, total)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return total


def render(prs, data, total, tmp, with_media):
    n = data["n"]
    timing = next((p for k, p in data["blocks"] if k == "timing"), "")
    blocks = [b for b in data["blocks"]
              if b[0] != "timing" and not (b[0] == "quote" and is_note(b[1]))]

    if n == 1:
        slide = title_slide(prs, data["title"], [
            "MemoActs 2026  ·  online intensive  ·  21.08",
            "Two hours. One thirty-second piece of film, four levels of intervention.",
        ])
        slide.notes_slide.notes_text_frame.text = notes_of(data)
        return

    picture = None
    if with_media:
        still, clip = media_for(data["raw"])
        if still:
            path = os.path.join(PROJECT, still)
            if os.path.exists(path):
                picture = path
            else:
                print("  missing still %s -- S%d has no picture" % (still, n))
        elif clip:
            picture = poster(clip, tmp)
    band = bool(picture) and aspect_of(picture) >= 2.2

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ground(slide, prs)
    width = FULL_WIDTH if (band or not picture) else SPLIT_WIDTH
    header(slide, n, timing)
    title(slide, data["title"], FULL_WIDTH if band else width)
    footer(slide, n, total)

    size = body_size(blocks, width)
    y = BODY_TOP
    for kind, payload in blocks:
        if kind == "para":
            cue = re.match(r"^\*\*(?:Screen|Show) `([^`]+)`\*\*", payload)
            if cue:
                y = draw_cue(slide, y, cue.group(1), width)
                rest = payload[cue.end():].lstrip(" —-").strip()
                if rest:
                    y = draw_para(slide, y, rest, width, size, muted=True)
                continue
            y = draw_para(slide, y, payload, width, size)
        elif kind == "h3":
            y = draw_para(slide, y + 0.05, "**%s**" % plain(payload), width, size + 2)
        elif kind in ("bullets", "numbers"):
            y = draw_list(slide, y, payload, width, size, kind == "numbers")
        elif kind == "table":
            y = draw_table(slide, y, payload, width, size)
        elif kind == "quote":
            draw = draw_callout if is_callout(payload) else draw_quote
            y = draw(slide, y, payload, width, size)

    if picture:
        place_picture(slide, picture, band, max(y + 0.10, BODY_TOP))

    slide.notes_slide.notes_text_frame.text = notes_of(data)


def run_sheet_slide(prs, preamble, total):
    """The contract, kept at the back where the presenter can find it."""
    rows = [table_row(l) for l in preamble if l.startswith("|")]
    rows = [r for r in rows if not is_separator_row(r)]
    if not rows:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ground(slide, prs)
    header(slide, total + 1, "")
    title(slide, "Presenter — the run sheet", FULL_WIDTH)
    draw_table(slide, BODY_TOP, rows, FULL_WIDTH, 14)
    footer(slide, total + 1, total + 1)
    slide.notes_slide.notes_text_frame.text = (
        "Not for the room. The contract: every block has a minute budget, and "
        "S13, the eight-node table, is the first thing to cut -- it is in the "
        "handout anyway."
    )


def main():
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--source", default=SOURCE)
    ap.add_argument("--out", default=DEFAULT_OUT)
    ap.add_argument("--no-media", action="store_true",
                    help="skip pictures -- text only, and no ffmpeg needed")
    args = ap.parse_args()

    if not os.path.exists(args.source):
        sys.exit("no run sheet at %s" % args.source)
    count = build(args.source, args.out, with_media=not args.no_media)
    print("%d slides + the run sheet -> %s" % (count, args.out))


if __name__ == "__main__":
    main()
